from __future__ import annotations

import os
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent


@dataclass(frozen=True)
class SlideSpec:
    number: int
    title: str
    bullets: list[str]
    subtitle: str | None = None


def _load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Carga una fuente razonable en Windows; cae a la default si falla."""

    # Windows suele tener Arial por defecto.
    candidates = []
    if bold:
        candidates.extend(
            [
                r"C:\\Windows\\Fonts\\arialbd.ttf",
                r"C:\\Windows\\Fonts\\calibrib.ttf",
            ]
        )
    candidates.extend(
        [
            r"C:\\Windows\\Fonts\\arial.ttf",
            r"C:\\Windows\\Fonts\\calibri.ttf",
            r"C:\\Windows\\Fonts\\segoeui.ttf",
        ]
    )

    for fp in candidates:
        try:
            if os.path.exists(fp):
                return ImageFont.truetype(fp, size=size)
        except Exception:
            continue

    # Fallback: fuente por defecto de PIL
    return ImageFont.load_default()


def _lerp(a: int, b: int, t: float) -> int:
    return int(round(a + (b - a) * t))


def _lerp_color(c1: tuple[int, int, int], c2: tuple[int, int, int], t: float) -> tuple[int, int, int]:
    return (_lerp(c1[0], c2[0], t), _lerp(c1[1], c2[1], t), _lerp(c1[2], c2[2], t))


def _draw_vertical_gradient(img: Image.Image, top: tuple[int, int, int], bottom: tuple[int, int, int]) -> None:
    draw = ImageDraw.Draw(img)
    w, h = img.size
    for y in range(h):
        t = y / max(h - 1, 1)
        draw.line([(0, y), (w, y)], fill=_lerp_color(top, bottom, t))


def _rounded_rect(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    radius: int,
    fill: tuple[int, int, int, int],
    outline: tuple[int, int, int, int] | None = None,
    width: int = 1,
) -> None:
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle([x1, y1, x2, y2], radius=radius, fill=fill, outline=outline, width=width)


def _shadow_card(
    base: Image.Image,
    xy: tuple[int, int, int, int],
    radius: int,
    fill: tuple[int, int, int, int],
    shadow: tuple[int, int, int, int] = (0, 0, 0, 110),
    offset: tuple[int, int] = (0, 10),
    blur: int = 16,
    outline: tuple[int, int, int, int] | None = None,
    outline_width: int = 1,
) -> None:
    """Dibuja una tarjeta con sombra suave (estilo corporativo)."""
    x1, y1, x2, y2 = xy

    shadow_layer = Image.new("RGBA", base.size, (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow_layer)
    sx1, sy1, sx2, sy2 = x1 + offset[0], y1 + offset[1], x2 + offset[0], y2 + offset[1]
    _rounded_rect(shadow_draw, (sx1, sy1, sx2, sy2), radius, fill=shadow)
    shadow_layer = shadow_layer.filter(ImageFilter.GaussianBlur(blur))
    base.alpha_composite(shadow_layer)

    card_layer = Image.new("RGBA", base.size, (0, 0, 0, 0))
    card_draw = ImageDraw.Draw(card_layer)
    _rounded_rect(card_draw, (x1, y1, x2, y2), radius, fill=fill, outline=outline, width=outline_width)
    base.alpha_composite(card_layer)


def _draw_arrow(
    draw: ImageDraw.ImageDraw,
    p1: tuple[int, int],
    p2: tuple[int, int],
    color: tuple[int, int, int, int],
    width: int = 4,
    head_len: int = 18,
    head_w: int = 10,
) -> None:
    x1, y1 = p1
    x2, y2 = p2
    draw.line([p1, p2], fill=color, width=width)

    ang = math.atan2(y2 - y1, x2 - x1)
    back = (x2 - head_len * math.cos(ang), y2 - head_len * math.sin(ang))
    left = (
        back[0] + head_w * math.cos(ang + math.pi / 2),
        back[1] + head_w * math.sin(ang + math.pi / 2),
    )
    right = (
        back[0] + head_w * math.cos(ang - math.pi / 2),
        back[1] + head_w * math.sin(ang - math.pi / 2),
    )
    draw.polygon([p2, left, right], fill=color)


def _text_center(draw: ImageDraw.ImageDraw, xy: tuple[int, int, int, int], text: str, font, fill) -> None:
    x1, y1, x2, y2 = xy
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text((x1 + (x2 - x1 - tw) / 2, y1 + (y2 - y1 - th) / 2), text, font=font, fill=fill)


def _draw_badge(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, font, fg, bg) -> None:
    pad_x, pad_y = 18, 10
    bbox = draw.textbbox((0, 0), text, font=font)
    w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
    rect = (x, y, x + w + 2 * pad_x, y + h + 2 * pad_y)
    draw.rounded_rectangle(rect, radius=20, fill=bg)
    draw.text((x + pad_x, y + pad_y), text, font=font, fill=fg)


def _draw_illustration(spec: SlideSpec, base: Image.Image, box: tuple[int, int, int, int]) -> None:
    """Ilustraciones simples tipo diagrama para cada slide (sin assets externos)."""
    draw = ImageDraw.Draw(base)
    x1, y1, x2, y2 = box

    # Paleta (RGBA)
    card = (17, 24, 39, 255)  # #111827
    border = (31, 41, 55, 255)  # #1F2937
    blue = (59, 130, 246, 255)  # #3B82F6
    cyan = (34, 211, 238, 255)  # #22D3EE
    green = (34, 197, 94, 255)  # #22C55E
    amber = (245, 158, 11, 255)  # #F59E0B
    text = (248, 250, 252, 255)  # #F8FAFC
    muted = (203, 213, 225, 255)  # #CBD5E1

    f_title = _load_font(28, bold=True)
    f_label = _load_font(22, bold=False)

    def node(r: tuple[int, int, int, int], label: str, accent_color: tuple[int, int, int, int]) -> None:
        _shadow_card(base, r, radius=18, fill=card, outline=border)
        # barra superior
        draw.rounded_rectangle((r[0] + 16, r[1] + 16, r[0] + 16 + 54, r[1] + 16 + 10), radius=8, fill=accent_color)
        _text_center(draw, (r[0] + 16, r[1] + 28, r[2] - 16, r[3] - 16), label, font=f_label, fill=text)

    if spec.number == 1:
        # Hero: 3 tarjetas representando Estrategia / Diseño / Desarrollo
        cols = ["Estrategia", "Diseño", "Desarrollo"]
        colors = [cyan, blue, green]
        gap = 22
        w = (x2 - x1 - 2 * gap) // 3
        h = 240
        top = y1 + 80
        for i, (lbl, c) in enumerate(zip(cols, colors, strict=False)):
            r = (x1 + i * (w + gap), top, x1 + i * (w + gap) + w, top + h)
            node(r, lbl, c)
        # Flecha suave indicando flujo
        mid_y = top + h + 60
        _draw_arrow(draw, (x1 + 40, mid_y), (x2 - 40, mid_y), color=(59, 130, 246, 200), width=6)
        _text_center(draw, (x1, mid_y + 14, x2, mid_y + 70), "De la idea a producción", font=f_title, fill=muted)

    elif spec.number == 2:
        # Checklist
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Brief", font=f_title, fill=text)
        items = [
            ("Objetivo", cyan),
            ("Audiencia", blue),
            ("Alcance", green),
            ("Métricas", amber),
        ]
        cy = y1 + 80
        for label, c in items:
            # check icon
            draw.ellipse((x1 + 28, cy, x1 + 28 + 26, cy + 26), fill=c)
            draw.line((x1 + 36, cy + 14, x1 + 40, cy + 18), fill=(0, 0, 0, 180), width=4)
            draw.line((x1 + 40, cy + 18, x1 + 50, cy + 8), fill=(0, 0, 0, 180), width=4)
            draw.text((x1 + 70, cy - 2), label, font=f_label, fill=muted)
            cy += 54

    elif spec.number == 3:
        # Arquitectura por capas
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Arquitectura", font=f_title, fill=text)
        margin = 26
        band_h = (y2 - y1 - 100) // 3
        labels = [("Front-end", cyan), ("Back-end", blue), ("Base de datos", green)]
        for i, (lbl, c) in enumerate(labels):
            r = (x1 + margin, y1 + 70 + i * (band_h + 12), x2 - margin, y1 + 70 + i * (band_h + 12) + band_h)
            _rounded_rect(draw, r, radius=18, fill=(255, 255, 255, 16), outline=(255, 255, 255, 26), width=1)
            draw.rounded_rectangle((r[0] + 14, r[1] + 14, r[0] + 14 + 54, r[1] + 14 + 10), radius=8, fill=c)
            _text_center(draw, (r[0], r[1] + 10, r[2], r[3]), lbl, font=f_label, fill=text)
            if i < 2:
                _draw_arrow(draw, (r[0] + 20, r[3] + 8), (r[0] + 20, r[3] + 8 + 24), color=(203, 213, 225, 180), width=4)

    elif spec.number == 4:
        # UI mockup + paleta
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "UI/UX", font=f_title, fill=text)
        # Browser window
        bw = int((x2 - x1) * 0.62)
        bh = int((y2 - y1) * 0.46)
        bx = x1 + 26
        by = y1 + 78
        _rounded_rect(draw, (bx, by, bx + bw, by + bh), radius=18, fill=(255, 255, 255, 18), outline=(255, 255, 255, 30))
        draw.rectangle((bx, by + 42, bx + bw, by + 44), fill=(255, 255, 255, 30))
        for i, c in enumerate([cyan, blue, green]):
            draw.ellipse((bx + 18 + i * 22, by + 14, bx + 18 + i * 22 + 12, by + 14 + 12), fill=c)
        # Phone mock
        px = bx + bw + 18
        py = by + 18
        pw = (x2 - x1) - (px - x1) - 26
        ph = bh + 80
        _rounded_rect(draw, (px, py, px + pw, py + ph), radius=28, fill=(255, 255, 255, 16), outline=(255, 255, 255, 30))
        draw.rounded_rectangle((px + 14, py + 20, px + pw - 14, py + 68), radius=14, fill=(255, 255, 255, 22))
        # Palette chips
        cy = by + bh + 28
        draw.text((x1 + 26, cy), "Paleta", font=f_label, fill=muted)
        chip_x = x1 + 26
        for c in [blue, cyan, green, amber]:
            draw.rounded_rectangle((chip_x, cy + 36, chip_x + 58, cy + 86), radius=14, fill=c)
            chip_x += 72

    elif spec.number == 5:
        # Layout/components
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Layout", font=f_title, fill=text)
        ux1, uy1, ux2, uy2 = x1 + 26, y1 + 78, x2 - 26, y2 - 26
        _rounded_rect(draw, (ux1, uy1, ux2, uy2), radius=18, fill=(255, 255, 255, 14), outline=(255, 255, 255, 28))
        # sections
        header = (ux1 + 18, uy1 + 18, ux2 - 18, uy1 + 84)
        hero = (ux1 + 18, uy1 + 100, ux2 - 18, uy1 + 220)
        cards = (ux1 + 18, uy1 + 238, ux2 - 18, uy2 - 18)
        for r, lbl, c in [(header, "Header", cyan), (hero, "Hero", blue), (cards, "Cards", green)]:
            _rounded_rect(draw, r, radius=14, fill=(255, 255, 255, 18), outline=(255, 255, 255, 30))
            draw.rounded_rectangle((r[0] + 12, r[1] + 12, r[0] + 12 + 50, r[1] + 12 + 10), radius=8, fill=c)
            draw.text((r[0] + 74, r[1] + 8), lbl, font=f_label, fill=muted)

    elif spec.number == 6:
        # Flujo interacción
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Interacción", font=f_title, fill=text)
        nodes = ["Usuario", "UI", "API", "Feedback"]
        colors = [cyan, blue, green, amber]
        cx = x1 + 38
        cy = y1 + 140
        w = 160
        h = 88
        gap = 26
        for i, (lbl, c) in enumerate(zip(nodes, colors, strict=False)):
            r = (cx + i * (w + gap), cy, cx + i * (w + gap) + w, cy + h)
            node(r, lbl, c)
            if i < len(nodes) - 1:
                _draw_arrow(draw, (r[2] + 8, cy + h // 2), (r[2] + gap - 8, cy + h // 2), color=(203, 213, 225, 180))

    elif spec.number == 7:
        # Request/response (pipeline)
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Backend", font=f_title, fill=text)
        steps = ["Route", "Controller", "Service", "Model", "DB"]
        colors = [cyan, blue, blue, green, green]
        left = x1 + 26
        top = y1 + 104
        box_h = 72
        gap = 18
        for i, (lbl, c) in enumerate(zip(steps, colors, strict=False)):
            r = (left, top + i * (box_h + gap), x2 - 26, top + i * (box_h + gap) + box_h)
            _rounded_rect(draw, r, radius=16, fill=(255, 255, 255, 16), outline=(255, 255, 255, 28))
            draw.rounded_rectangle((r[0] + 16, r[1] + 16, r[0] + 16 + 54, r[1] + 16 + 10), radius=8, fill=c)
            draw.text((r[0] + 86, r[1] + 18), lbl, font=f_label, fill=muted)
            if i < len(steps) - 1:
                _draw_arrow(draw, (r[0] + 34, r[3] + 6), (r[0] + 34, r[3] + 6 + 18), color=(203, 213, 225, 180), width=4)

    elif spec.number == 8:
        # Tablas + relaciones
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Datos", font=f_title, fill=text)
        tables = [
            ("users", (x1 + 40, y1 + 110, x1 + 260, y1 + 250), blue),
            ("leads", (x1 + 320, y1 + 190, x1 + 540, y1 + 330), cyan),
            ("projects", (x1 + 180, y1 + 360, x1 + 420, y1 + 500), green),
        ]
        centers: dict[str, tuple[int, int]] = {}
        for name, r, c in tables:
            _rounded_rect(draw, r, radius=16, fill=(255, 255, 255, 16), outline=(255, 255, 255, 28))
            draw.rectangle((r[0], r[1] + 44, r[2], r[1] + 46), fill=(255, 255, 255, 30))
            draw.rounded_rectangle((r[0] + 14, r[1] + 16, r[0] + 14 + 54, r[1] + 16 + 10), radius=8, fill=c)
            draw.text((r[0] + 78, r[1] + 10), name, font=f_label, fill=muted)
            centers[name] = ((r[0] + r[2]) // 2, (r[1] + r[3]) // 2)
        # relations
        _draw_arrow(draw, centers["users"], centers["leads"], color=(203, 213, 225, 180), width=4)
        _draw_arrow(draw, centers["users"], centers["projects"], color=(203, 213, 225, 180), width=4)

    elif spec.number == 9:
        # Gauge + checklist QA
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "QA", font=f_title, fill=text)
        # Gauge
        gx, gy = x1 + 110, y1 + 170
        r = 120
        draw.arc((gx - r, gy - r, gx + r, gy + r), start=200, end=340, fill=(203, 213, 225, 120), width=16)
        draw.arc((gx - r, gy - r, gx + r, gy + r), start=200, end=300, fill=green, width=16)
        _text_center(draw, (gx - 60, gy - 18, gx + 60, gy + 40), "A+", font=_load_font(48, bold=True), fill=text)
        # Items
        items = [("Tests", blue), ("Performance", cyan), ("Seguridad", amber)]
        iy = y1 + 120
        for i, (lbl, c) in enumerate(items):
            rx1 = x1 + 270
            ry1 = iy + i * 96
            rr = (rx1, ry1, x2 - 26, ry1 + 70)
            node(rr, lbl, c)

    elif spec.number == 10:
        # Pipeline CI/CD
        _shadow_card(base, (x1, y1, x2, y2), radius=22, fill=card, outline=border)
        draw.text((x1 + 26, y1 + 24), "Deploy", font=f_title, fill=text)
        steps = ["Git", "Build", "Test", "Deploy", "Monitor"]
        colors = [cyan, blue, blue, green, amber]
        gap = 18
        w = (x2 - x1 - 2 * 26 - (len(steps) - 1) * gap) // len(steps)
        y = y1 + 190
        for i, (lbl, c) in enumerate(zip(steps, colors, strict=False)):
            rx1 = x1 + 26 + i * (w + gap)
            rr = (rx1, y, rx1 + w, y + 86)
            node(rr, lbl, c)
            if i < len(steps) - 1:
                _draw_arrow(draw, (rr[2] + 8, y + 43), (rr[2] + gap - 8, y + 43), color=(203, 213, 225, 180))


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]

    lines: list[str] = []
    current = words[0]
    for w in words[1:]:
        candidate = f"{current} {w}"
        bbox = draw.textbbox((0, 0), candidate, font=font)
        if bbox[2] <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = w
    lines.append(current)
    return lines


def _render_slide_png(spec: SlideSpec, out_path: Path) -> None:
    # Formato 16:9 (Full HD)
    W, H = 1920, 1080

    # Paleta (RGBA)
    bg = (11, 18, 32, 255)  # #0B1220
    card = (17, 24, 39, 255)  # #111827
    border = (31, 41, 55, 255)  # #1F2937
    blue = (59, 130, 246, 255)  # #3B82F6
    cyan = (34, 211, 238, 255)  # #22D3EE
    text_main = (248, 250, 252, 255)  # #F8FAFC
    text_muted = (203, 213, 225, 255)  # #CBD5E1

    base = Image.new("RGBA", (W, H), color=bg)

    # Decoración abstracta (circulos con alpha)
    deco = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(deco)
    d.ellipse((W - 760, -260, W + 260, 760), fill=(59, 130, 246, 70))
    d.ellipse((-420, H - 620, 620, H + 420), fill=(34, 211, 238, 60))
    d.ellipse((W - 520, H - 520, W + 140, H + 140), fill=(34, 197, 94, 42))
    base.alpha_composite(deco.filter(ImageFilter.GaussianBlur(8)))

    draw = ImageDraw.Draw(base)

    # Header corporativo
    header_h = 92
    _rounded_rect(draw, (0, 0, W, header_h), radius=0, fill=(8, 10, 18, 220))
    draw.line((0, header_h, W, header_h), fill=(255, 255, 255, 18), width=2)
    # Barra izquierda (accent)
    accent_rgb = Image.new("RGB", (16, H), (0, 0, 0))
    _draw_vertical_gradient(accent_rgb, (59, 130, 246), (34, 211, 238))
    base.alpha_composite(accent_rgb.convert("RGBA"), (0, 0))

    # Fuentes
    title_font = _load_font(60, bold=True)
    subtitle_font = _load_font(30, bold=False)
    bullet_font = _load_font(34, bold=False)
    small_font = _load_font(22, bold=False)
    header_font = _load_font(22, bold=True)

    # Header contents
    draw.text((34, 28), f"{spec.number:02d}", font=_load_font(28, bold=True), fill=text_main)
    _draw_badge(draw, 92, 22, "Proceso Web", font=_load_font(20, bold=True), fg=text_main, bg=(59, 130, 246, 160))

    # Marca derecha
    brand = "GusGus Web"
    brand_bbox = draw.textbbox((0, 0), brand, font=header_font)
    draw.text((W - 40 - (brand_bbox[2] - brand_bbox[0]), 28), brand, font=header_font, fill=text_muted)

    # Layout: texto a la izquierda + diagrama a la derecha
    pad_x = 90
    content_top = header_h + 60
    content_bottom = H - 86
    left_w = int(W * 0.56)
    gap = 44
    right_x1 = pad_x + left_w + gap

    if spec.number == 1:
        # Portada: centrado + hero cards
        title_lines = _wrap_text(draw, spec.title, _load_font(76, bold=True), W - 2 * pad_x)
        tfont = _load_font(76, bold=True)
        y = 260
        for line in title_lines[:2]:
            _text_center(draw, (pad_x, y, W - pad_x, y + 90), line, font=tfont, fill=text_main)
            y += 92
        if spec.subtitle:
            _text_center(draw, (pad_x, y + 10, W - pad_x, y + 60), spec.subtitle, font=_load_font(34, bold=False), fill=text_muted)

        # Ilustración ocupa el centro inferior
        ill_box = (pad_x + 140, 520, W - pad_x - 140, H - 150)
        _draw_illustration(spec, base, ill_box)

        # Footer
        draw.text((pad_x, H - 58), "Presentación corporativa", font=small_font, fill=text_muted)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        base.convert("RGB").save(out_path, format="PNG", optimize=True)
        return

    # Tarjeta texto
    text_card = (pad_x, content_top, pad_x + left_w, content_bottom)
    _shadow_card(base, text_card, radius=28, fill=card, outline=border)

    # Diagrama
    ill_card = (right_x1, content_top, W - pad_x, content_bottom)
    _shadow_card(base, ill_card, radius=28, fill=card, outline=border)
    _draw_illustration(spec, base, (ill_card[0] + 18, ill_card[1] + 18, ill_card[2] - 18, ill_card[3] - 18))

    # Título dentro de tarjeta de texto
    title_x = text_card[0] + 34
    y = text_card[1] + 34
    max_title_w = text_card[2] - title_x - 34
    title_lines = _wrap_text(draw, spec.title, title_font, max_title_w)
    if len(title_lines) > 2:
        title_font = _load_font(54, bold=True)
        title_lines = _wrap_text(draw, spec.title, title_font, max_title_w)
    for line in title_lines[:2]:
        draw.text((title_x, y), line, font=title_font, fill=text_main)
        y += int((title_font.size or 60) * 1.12)

    if spec.subtitle:
        draw.text((title_x, y + 6), spec.subtitle, font=subtitle_font, fill=text_muted)
        y += 64
    else:
        y += 34

    # Bullets
    bullet_area_w = text_card[2] - title_x - 34
    bullet_indent = 42
    bullet_gap = 18

    for b in spec.bullets:
        draw.text((title_x, y), "•", font=bullet_font, fill=(34, 211, 238, 255))
        wrapped = _wrap_text(draw, b, bullet_font, bullet_area_w - bullet_indent)
        for i, line in enumerate(wrapped):
            draw.text((title_x + bullet_indent, y), line, font=bullet_font, fill=text_main)
            if i < len(wrapped) - 1:
                y += int((bullet_font.size or 34) * 1.18)
        y += int((bullet_font.size or 34) * 1.18) + bullet_gap
        if y > text_card[3] - 120:
            break

    # Footer
    footer_text_left = "Proceso de creación de una página web"
    footer_text_right = "Laravel • Blade • Tailwind • Vite"
    footer_y = H - 58
    draw.text((pad_x, footer_y), footer_text_left, font=small_font, fill=text_muted)
    bbox = draw.textbbox((0, 0), footer_text_right, font=small_font)
    draw.text((W - pad_x - (bbox[2] - bbox[0]), footer_y), footer_text_right, font=small_font, fill=text_muted)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    base.convert("RGB").save(out_path, format="PNG", optimize=True)


def _build_pptx(image_paths: Iterable[Path], out_pptx: Path) -> None:
    prs = Presentation()
    # 16:9 (widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def main() -> int:
    slides: list[SlideSpec] = [
        SlideSpec(
            1,
            "Proceso de creación de una página web",
            [
                "De la idea al despliegue: pasos, roles y entregables.",
                "Ejemplo de stack: Laravel + Blade + Tailwind + Vite.",
            ],
            subtitle="Guía práctica (resumen)",
        ),
        SlideSpec(
            2,
            "1) Brief y objetivos",
            [
                "Definir propósito (ventas, leads, portafolio, soporte).",
                "Identificar audiencia, tono, y propuesta de valor.",
                "Listar páginas/funcionalidades y métricas de éxito.",
            ],
        ),
        SlideSpec(
            3,
            "2) Planificación y arquitectura",
            [
                "Mapa del sitio (secciones) y flujo del usuario.",
                "Decidir stack: Front (UI), Back (lógica) y DB (datos).",
                "Estructurar rutas, componentes, servicios y permisos.",
            ],
        ),
        SlideSpec(
            4,
            "3) Diseño UI/UX",
            [
                "Wireframes → prototipo (Figma) → diseño final.",
                "Guía visual: colores, tipografía, espaciado, iconografía.",
                "Diseño responsive (mobile-first) y accesibilidad.",
            ],
        ),
        SlideSpec(
            5,
            "4) Maquetación (Front-end)",
            [
                "HTML semántico + CSS (Tailwind) para construir la interfaz.",
                "Componentes reutilizables (Blade components) y layouts.",
                "Optimización de assets (Vite): CSS/JS minificados.",
            ],
        ),
        SlideSpec(
            6,
            "5) Interactividad y experiencia",
            [
                "JavaScript para formularios, menús, validaciones y UX.",
                "Estados de carga, errores y mensajes claros al usuario.",
                "Accesibilidad: foco, ARIA, contraste, teclado.",
            ],
        ),
        SlideSpec(
            7,
            "6) Back-end y reglas de negocio",
            [
                "Rutas → Controladores → Servicios → Modelos.",
                "Validación, autenticación/autorización y seguridad.",
                "Integraciones: email, APIs, almacenamiento de archivos.",
            ],
        ),
        SlideSpec(
            8,
            "7) Base de datos",
            [
                "Modelo de datos (tablas y relaciones).",
                "Migraciones para versionar cambios y mantener consistencia.",
                "Seeders/Factories para datos de prueba y entornos.",
            ],
        ),
        SlideSpec(
            9,
            "8) Pruebas, rendimiento y calidad",
            [
                "Pruebas funcionales (flujos), validación y casos borde.",
                "Performance: caché, queries eficientes, imágenes optimizadas.",
                "Revisión de seguridad: CSRF, XSS, permisos, logs.",
            ],
        ),
        SlideSpec(
            10,
            "9) Despliegue y mantenimiento",
            [
                "Deploy: variables de entorno, build de assets y migraciones.",
                "Monitoreo, backups, métricas y corrección de errores.",
                "Iteración continua: contenido, SEO, mejoras UX.",
            ],
        ),
    ]

    image_paths: list[Path] = []
    for spec in slides:
        out = ROOT / f"slide-{spec.number:02d}.png"
        _render_slide_png(spec, out)
        image_paths.append(out)

    out_pptx = ROOT / "Proceso_creacion_pagina_web.pptx"
    _build_pptx(image_paths, out_pptx)

    print(f"OK: imágenes generadas en: {ROOT}")
    print(f"OK: pptx generado: {out_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

